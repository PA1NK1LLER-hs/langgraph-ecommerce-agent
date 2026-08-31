"""多格式文档解析器。

支持的格式：
- PDF: unstructured + pdfplumber（表格提取）
- DOCX: python-docx
- XLSX: openpyxl → Markdown 表格
- PPTX: python-pptx → 幻灯片文本
- HTML/URL: trafilatura（正文提取）+ BeautifulSoup（元数据）
- 图片: base64 编码供 Vision LLM 处理
- 纯文本: 直接返回

分块策略：
- semantic: 按段落 + 标题边界（保留文档结构）
- fixed: 固定 512 tokens + 128 overlap
- recursive: 按 \\n\\n → \\n → 。 → ， 递归分割
"""

from __future__ import annotations

import io
import logging
import mimetypes
import os
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# 数据结构
# ---------------------------------------------------------------------------


@dataclass
class Chunk:
    """分块结果。"""

    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    chunk_index: int = 0
    chunk_type: str = "text"

    def to_dict(self) -> dict[str, Any]:
        return {
            "text": self.text,
            "metadata": self.metadata,
            "chunk_index": self.chunk_index,
            "chunk_type": self.chunk_type,
        }


@dataclass
class ParsedDocument:
    """解析后的文档，包含原始文本和元数据。"""

    text: str
    metadata: dict[str, Any] = field(default_factory=dict)
    pages: list[str] = field(default_factory=list)
    tables: list[str] = field(default_factory=list)

    def to_chunks(
        self,
        strategy: str = "semantic",
        chunk_size: int = 512,
        chunk_overlap: int = 128,
    ) -> list[Chunk]:
        """将文档文本按策略分块。"""
        chunker = _get_chunker(strategy, chunk_size, chunk_overlap)
        chunks = chunker.split(self.text)
        result = []
        for i, ch in enumerate(chunks):
            result.append(Chunk(
                text=ch,
                metadata={
                    **self.metadata,
                    "chunk_of": len(chunks),
                },
                chunk_index=i,
            ))
        return result


# ---------------------------------------------------------------------------
# 分块器
# ---------------------------------------------------------------------------


class _SemanticChunker:
    """按段落 + 标题边界分块，保留文档结构。"""

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 128):
        self._size = chunk_size
        self._overlap = chunk_overlap

    def split(self, text: str) -> list[str]:
        # 先按双换行分段
        paragraphs = re.split(r"\n\s*\n", text)
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0

        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            # 估算 token 数（中文字符 ≈ 1.5 tokens，英文单词 ≈ 1.3 tokens）
            para_tokens = self._estimate_tokens(para)

            if current_len + para_tokens > self._size and current:
                chunks.append("\n\n".join(current))
                # overlap: 保留最后一个段落
                if len(current) > 1 and self._overlap > 0:
                    current = [current[-1]]
                    current_len = self._estimate_tokens(current[-1])
                else:
                    current = []
                    current_len = 0

            current.append(para)
            current_len += para_tokens

        if current:
            chunks.append("\n\n".join(current))

        return chunks if chunks else [text]

    @staticmethod
    def _estimate_tokens(text: str) -> int:
        """粗略估算 token 数。"""
        chinese_chars = len(re.findall(r"[一-鿿]", text))
        other_chars = len(text) - chinese_chars
        return int(chinese_chars * 1.5 + other_chars * 0.3)


class _FixedChunker:
    """固定大小分块。"""

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 128):
        self._size = chunk_size
        self._overlap = chunk_overlap

    def split(self, text: str) -> list[str]:
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + self._size, len(text))
            chunks.append(text[start:end])
            start += self._size - self._overlap
        return chunks if chunks else [text]


class _RecursiveChunker:
    """递归分割：按 \\n\\n → \\n → 。 → ， 逐级尝试。"""

    def __init__(self, chunk_size: int = 512, chunk_overlap: int = 128):
        self._size = chunk_size
        self._overlap = chunk_overlap

    def split(self, text: str) -> list[str]:
        separators = ["\n\n", "\n", "。", "，", ". ", ", ", " "]
        return self._split_recursive(text, separators)

    def _split_recursive(self, text: str, separators: list[str]) -> list[str]:
        if not separators:
            # 最终回退：按固定长度强制分割
            return self._force_split(text)

        sep = separators[0]
        remaining = separators[1:]

        parts = text.split(sep)
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0

        for part in parts:
            part_len = len(part)
            if current_len + part_len > self._size and current:
                chunks.append(sep.join(current))
                if len(current) > 1:
                    current = [current[-1]]
                    current_len = len(current[-1])
                else:
                    current = []
                    current_len = 0

            if part_len > self._size and remaining:
                sub_chunks = self._split_recursive(part, remaining)
                for sc in sub_chunks:
                    if current_len + len(sc) > self._size and current:
                        chunks.append(sep.join(current))
                        current = []
                        current_len = 0
                    current.append(sc)
                    current_len += len(sc)
            else:
                current.append(part)
                current_len += part_len

        if current:
            chunks.append(sep.join(current))

        # 如果只有一个 chunk 且大于 chunk_size，强制分割
        if len(chunks) == 1 and len(chunks[0]) > self._size:
            return self._force_split(chunks[0])

        return chunks if chunks else [text]

    def _force_split(self, text: str) -> list[str]:
        """强制按固定长度分割（无分隔符时的最终回退）。"""
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + self._size, len(text))
            chunks.append(text[start:end])
            start += self._size - self._overlap
        return chunks if chunks else [text]


_CHUNKERS = {
    "semantic": _SemanticChunker,
    "fixed": _FixedChunker,
    "recursive": _RecursiveChunker,
}


def _get_chunker(strategy: str, chunk_size: int, chunk_overlap: int):
    cls = _CHUNKERS.get(strategy, _SemanticChunker)
    return cls(chunk_size, chunk_overlap)


# ---------------------------------------------------------------------------
# 主解析器
# ---------------------------------------------------------------------------


class DocumentParser:
    """多格式文档解析器，自动检测文件类型并分发到对应解析器。"""

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".doc", ".xlsx", ".xls",
        ".pptx", ".ppt", ".txt", ".md", ".csv",
        ".html", ".htm", ".json", ".xml",
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp",
    }

    def parse_file(self, file_path: str | Path) -> ParsedDocument:
        """解析文件，自动检测格式。

        Args:
            file_path: 文件路径。

        Returns:
            ParsedDocument 包含文本、元数据、表格等。

        Raises:
            ValueError: 不支持的文件格式。
            FileNotFoundError: 文件不存在。
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        ext = path.suffix.lower()
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"不支持的文件格式: {ext}（支持的格式: {sorted(self.SUPPORTED_EXTENSIONS)}）")

        metadata = {
            "source": str(path),
            "filename": path.name,
            "extension": ext,
            "size_bytes": path.stat().st_size,
        }

        mime_type, _ = mimetypes.guess_type(str(path))

        if ext == ".pdf":
            return self._parse_pdf(path, metadata)
        elif ext in (".docx", ".doc"):
            return self._parse_docx(path, metadata)
        elif ext in (".xlsx", ".xls"):
            return self._parse_xlsx(path, metadata)
        elif ext in (".pptx", ".ppt"):
            return self._parse_pptx(path, metadata)
        elif ext in (".html", ".htm"):
            return self._parse_html(path, metadata)
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            return self._parse_image(path, metadata)
        else:
            return self._parse_text(path, metadata)

    def parse_url(self, url: str) -> ParsedDocument:
        """从 URL 抓取并解析网页。

        Args:
            url: 网页 URL。

        Returns:
            ParsedDocument 包含正文文本和元数据。
        """
        metadata = {"source": url, "filename": url.rsplit("/", 1)[-1] or "index.html"}

        text, meta = self._fetch_url(url)
        metadata.update(meta)

        return ParsedDocument(
            text=text,
            metadata=metadata,
        )

    # ------------------------------------------------------------------
    # 内部解析器
    # ------------------------------------------------------------------

    @staticmethod
    def _parse_text(path: Path, metadata: dict) -> ParsedDocument:
        """解析纯文本文件。"""
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = path.read_text(encoding="gbk", errors="replace")
        return ParsedDocument(text=text, metadata=metadata)

    @staticmethod
    def _parse_html(path: Path, metadata: dict) -> ParsedDocument:
        """解析 HTML 文件，提取正文。"""
        try:
            import trafilatura
        except ImportError:
            logger.warning("trafilatura 未安装，回退到纯文本读取")
            return DocumentParser._parse_text(path, metadata)

        html = path.read_text(encoding="utf-8", errors="replace")
        text = trafilatura.extract(html, include_comments=False, include_tables=True) or ""
        if not text:
            # 回退：BeautifulSoup 提取所有文本
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html, "html.parser")
                text = soup.get_text(separator="\n", strip=True)
            except ImportError:
                text = html

        # 提取 title
        title = ""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(html, "html.parser")
            title_tag = soup.find("title")
            if title_tag:
                title = title_tag.get_text(strip=True)
                metadata["title"] = title
        except Exception:
            pass

        return ParsedDocument(text=text, metadata=metadata)

    @staticmethod
    def _parse_pdf(path: Path, metadata: dict) -> ParsedDocument:
        """解析 PDF 文件。"""
        text_parts: list[str] = []
        tables: list[str] = []

        # 优先使用 pdfplumber 提取文本和表格
        try:
            import pdfplumber
            with pdfplumber.open(str(path)) as pdf:
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
                    # 提取表格
                    page_tables = page.extract_tables()
                    for j, tbl in enumerate(page_tables):
                        if tbl:
                            md_table = _table_to_markdown(tbl)
                            tables.append(md_table)
                            text_parts.append(f"[表格 {i+1}-{j+1}]\n{md_table}")
        except ImportError:
            logger.info("pdfplumber 未安装，使用 unstructured 解析 PDF")
        except Exception as exc:
            logger.warning("pdfplumber 解析 PDF 失败: %s，尝试 unstructured", exc)

        # 如果 pdfplumber 没提取到文本，回退到 unstructured
        if not text_parts:
            try:
                import unstructured.partition.pdf as updf
                elements = updf.partition_pdf(str(path))
                text_parts = [str(el) for el in elements if str(el).strip()]
            except ImportError:
                logger.warning("unstructured 未安装，尝试 PyPDF2")
            except Exception as exc:
                logger.warning("unstructured 解析 PDF 失败: %s，尝试 PyPDF2", exc)

        # 最终回退到 PyPDF2
        if not text_parts:
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(str(path))
                metadata.setdefault("total_pages", len(reader.pages))
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            except ImportError:
                raise RuntimeError("无法解析 PDF：请安装 pdfplumber、unstructured 或 PyPDF2")
            except Exception as exc:
                raise RuntimeError(f"PDF 解析失败: {exc}")

        return ParsedDocument(
            text="\n\n".join(text_parts),
            metadata=metadata,
            tables=tables,
        )

    @staticmethod
    def _parse_docx(path: Path, metadata: dict) -> ParsedDocument:
        """解析 DOCX 文件。"""
        try:
            import docx
            doc = docx.Document(str(path))

            # 提取段落
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

            # 提取表格
            tables = []
            for i, tbl in enumerate(doc.tables):
                rows = []
                for row in tbl.rows:
                    cells = [cell.text for cell in row.cells]
                    rows.append(cells)
                md_table = _table_to_markdown(rows)
                tables.append(md_table)
                paragraphs.append(f"[表格 {i+1}]\n{md_table}")

            text = "\n\n".join(paragraphs)
            return ParsedDocument(text=text, metadata=metadata, tables=tables)
        except ImportError:
            raise RuntimeError("解析 DOCX 需要安装 python-docx")

    @staticmethod
    def _parse_xlsx(path: Path, metadata: dict) -> ParsedDocument:
        """解析 XLSX/CSV 文件，转换为 Markdown 表格。"""
        ext = path.suffix.lower()
        try:
            import pandas as pd
            if ext == ".csv":
                df = pd.read_csv(str(path))
            else:
                # XLSX/XLS - 读取所有 sheet
                xlsx = pd.ExcelFile(str(path))
                tables = []
                for sheet_name in xlsx.sheet_names:
                    df = pd.read_excel(xlsx, sheet_name=sheet_name)
                    if not df.empty:
                        df = df.fillna("")
                        # 限制行数避免过长
                        if len(df) > 100:
                            df = df.head(100)
                        md = df.to_markdown(index=False) if hasattr(df, "to_markdown") else df.to_string(index=False)
                        tables.append(f"## Sheet: {sheet_name}\n{md}")
                text = "\n\n".join(tables)
                return ParsedDocument(
                    text=text,
                    metadata={**metadata, "sheets": xlsx.sheet_names},
                    tables=tables,
                )

            df = df.fillna("")
            if len(df) > 100:
                df = df.head(100)
            md = df.to_markdown(index=False) if hasattr(df, "to_markdown") else df.to_string(index=False)
            return ParsedDocument(text=md, metadata=metadata)
        except ImportError:
            raise RuntimeError("解析表格文件需要安装 pandas 和 openpyxl（.xlsx）或 tabulate")

    @staticmethod
    def _parse_pptx(path: Path, metadata: dict) -> ParsedDocument:
        """解析 PPTX 文件，提取每页文本。"""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_parts = [f"## 幻灯片 {i+1}"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_parts.append(para.text.strip())
                slides_text.append("\n".join(slide_parts))
            text = "\n\n".join(slides_text)
            return ParsedDocument(text=text, metadata=metadata)
        except ImportError:
            raise RuntimeError("解析 PPTX 需要安装 python-pptx")

    @staticmethod
    def _parse_image(path: Path, metadata: dict) -> ParsedDocument:
        """解析图片 — 转为 base64 供 Vision LLM 处理。

        实际 OCR/描述由 Vision LLM 完成（见 agent.vision.describe_image）。
        这里只做 base64 编码并存到 metadata.image_base64；text 放占位符
        ``[IMAGE: name]``，避免把 base64 大 blob 塞进分块文本被向量化/索引
        （base64 对检索毫无意义，还会撑爆 chunk）。入库前由
        agent.vision.enrich_image_chunks 用视觉描述替换占位符。
        """
        import base64
        mime_type, _ = mimetypes.guess_type(str(path))
        if not mime_type:
            mime_type = "image/png"

        with open(path, "rb") as f:
            data = base64.b64encode(f.read()).decode("ascii")

        data_uri = f"data:{mime_type};base64,{data}"
        text = f"[IMAGE: {path.name}]"
        return ParsedDocument(text=text, metadata={**metadata, "image_base64": data_uri, "mime_type": mime_type})

    @staticmethod
    def _fetch_url(url: str) -> tuple[str, dict]:
        """抓取网页并提取正文。"""
        import urllib.request
        import urllib.error

        meta: dict[str, str] = {}

        try:
            req = urllib.request.Request(url, headers={
                "User-Agent": "Mozilla/5.0 (compatible; LangGraph-Agent/1.0)",
            })
            with urllib.request.urlopen(req, timeout=30) as resp:
                html = resp.read().decode("utf-8", errors="replace")
                meta["final_url"] = resp.geturl()
                meta["status_code"] = str(resp.getcode())
        except urllib.error.HTTPError as e:
            raise RuntimeError(f"HTTP {e.code}: 无法获取 {url}")
        except Exception as e:
            raise RuntimeError(f"获取 {url} 失败: {e}")

        # 正文提取
        try:
            import trafilatura
            text = trafilatura.extract(html, include_comments=False, include_tables=True) or ""
        except ImportError:
            text = ""

        if not text:
            # 回退到 BeautifulSoup
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(html, "html.parser")
                # 移除 script/style/nav/footer
                for tag in soup(["script", "style", "nav", "footer", "header"]):
                    tag.decompose()
                text = soup.get_text(separator="\n", strip=True)
                title_tag = soup.find("title")
                if title_tag:
                    meta["title"] = title_tag.get_text(strip=True)
            except ImportError:
                text = html

        return text, meta


# ---------------------------------------------------------------------------
# 辅助函数
# ---------------------------------------------------------------------------


def _table_to_markdown(rows: list[list[str | None]]) -> str:
    """将二维列表转为 Markdown 表格。"""
    if not rows:
        return ""

    # 清理 None 值
    clean = [[str(cell) if cell is not None else "" for cell in row] for row in rows]
    max_cols = max(len(row) for row in clean)
    # 补齐所有行到相同列数
    for row in clean:
        while len(row) < max_cols:
            row.append("")

    lines = []
    # header
    header = clean[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join(["---"] * max_cols) + " |")
    # body
    for row in clean[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# 模块级便捷函数
# ---------------------------------------------------------------------------


_parser: DocumentParser | None = None


def get_parser() -> DocumentParser:
    global _parser
    if _parser is None:
        _parser = DocumentParser()
    return _parser


def parse_file(file_path: str, chunk_strategy: str = "semantic") -> tuple[ParsedDocument, list[Chunk]]:
    """解析文件并分块。

    Returns:
        (解析的文档, 分块列表)
    """
    doc = get_parser().parse_file(file_path)
    chunks = doc.to_chunks(strategy=chunk_strategy)
    return doc, chunks


def parse_url(url: str, chunk_strategy: str = "semantic") -> tuple[ParsedDocument, list[Chunk]]:
    """抓取 URL 并分块。"""
    doc = get_parser().parse_url(url)
    chunks = doc.to_chunks(strategy=chunk_strategy)
    return doc, chunks
